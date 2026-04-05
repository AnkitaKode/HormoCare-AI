from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PINK = RGBColor(0xFF, 0x6B, 0x9D)
DARK_PINK = RGBColor(0xC4, 0x45, 0x69)
DARK = RGBColor(0x2C, 0x3E, 0x50)
GRAY = RGBColor(0x7F, 0x8C, 0x8D)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF8, 0xF9, 0xFA)
GREEN = RGBColor(0x2E, 0xCC, 0x71)
RED = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
TEAL = RGBColor(0x4E, 0xCD, 0xC4)
BLUE = RGBColor(0x45, 0xB7, 0xD1)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, r, g, b):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def add_gradient_bar(slide, top=0, height=Inches(0.08)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, top, W, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = PINK
    shape.line.fill.background()


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", line_spacing=1.3):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.line_spacing = Pt(font_size * line_spacing)
    return txbox


def add_rich_text_box(slide, left, top, width, height, lines,
                      font_name="Calibri"):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line_data.get("text", "")
        p.font.size = Pt(line_data.get("size", 18))
        p.font.bold = line_data.get("bold", False)
        p.font.color.rgb = line_data.get("color", DARK)
        p.font.name = font_name
        p.alignment = line_data.get("align", PP_ALIGN.LEFT)
        p.space_after = Pt(line_data.get("space_after", 6))
        p.space_before = Pt(line_data.get("space_before", 0))
    return txbox


def add_rounded_card(slide, left, top, width, height,
                     fill_rgb=WHITE, border_rgb=None):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if border_rgb:
        shape.line.color.rgb = border_rgb
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_circle(slide, left, top, size, fill_rgb):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    shape.line.fill.background()
    return shape


def add_accent_line(slide, left, top, width):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.06)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PINK
    shape.line.fill.background()


# ─── SLIDE 1: Title Slide ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide, top=0)
add_gradient_bar(slide, top=H - Inches(0.08))

add_circle(slide, Inches(10.5), Inches(1), Inches(4), RGBColor(0xFF, 0xF0, 0xF5))
add_circle(slide, Inches(-1), Inches(4), Inches(3), RGBColor(0xFF, 0xF0, 0xF5))

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1),
             "🌸", font_size=60, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             "HormoCare AI", font_size=54, bold=True, color=PINK,
             alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(3.7), Inches(9), Inches(0.8),
             "AI-Powered Early PCOS Detection for Indian Women",
             font_size=26, color=DARK, alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(4.6), Inches(2.3))

add_text_box(slide, Inches(2), Inches(5.0), Inches(9), Inches(1.2),
             "Because no woman should suffer in silence\n"
             "from an undiagnosed condition that affects 1 in 5.",
             font_size=18, color=GRAY, alignment=PP_ALIGN.CENTER,
             line_spacing=1.5)


# ─── SLIDE 2: The Problem ──────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "The Silent Health Crisis", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.6),
             "PCOS is the most common hormonal disorder in women of reproductive age — yet it remains invisible.",
             font_size=16, color=GRAY)

stats = [
    ("20%", "of Indian women\nhave PCOS", "100+ million women affected", RED),
    ("70%", "of cases go\nundiagnosed", "Average delay: 2-3 years", ORANGE),
    ("50%", "risk of developing\nType 2 Diabetes", "Without early intervention", DARK_PINK),
    ("2–3 yrs", "average time\nto diagnosis", "Precious time lost", TEAL),
]

for i, (num, label, detail, accent) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.2)
    add_rounded_card(slide, x, y, Inches(2.8), Inches(2.7), LIGHT_BG)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, y, Inches(2.8), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.8),
                 num, font_size=36, bold=True, color=accent,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.2), Inches(2.2), Inches(0.8),
                 label, font_size=14, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(2.0), Inches(2.2), Inches(0.5),
                 detail, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

consequences = [
    "⚠️  Irregular periods are normalised and ignored",
    "⚠️  Women don't track symptoms systematically",
    "⚠️  Delayed diagnosis → infertility, obesity, diabetes, depression",
    "⚠️  No early-screening tool designed for Indian women",
]
for j, txt in enumerate(consequences):
    add_text_box(slide, Inches(0.8), Inches(5.2 + j * 0.42), Inches(11), Inches(0.4),
                 txt, font_size=14, color=DARK)


# ─── SLIDE 3: Why Current Solutions Fail ───────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Why Current Solutions Fail", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

failures = [
    ("Generic Period Trackers",
     "Only track dates — completely miss hormonal symptoms like acne, hair loss, weight gain, and mood changes."),
    ("No Risk Assessment",
     "Zero intelligence — they log data but never analyse patterns or warn about PCOS risk."),
    ("Not India-Specific",
     "Western-centric apps ignore Indian dietary habits, genetic predisposition, and lifestyle factors."),
    ("Reactive, Not Preventive",
     "Women only seek help after complications appear. There is no tool pushing early detection."),
]

for i, (title, desc) in enumerate(failures):
    y = Inches(1.6 + i * 1.3)
    add_rounded_card(slide, Inches(0.8), y, Inches(5.5), Inches(1.1),
                     RGBColor(0xFF, 0xF5, 0xF5), RED)
    add_text_box(slide, Inches(1.1), y + Inches(0.1), Inches(5), Inches(0.4),
                 "✗  " + title, font_size=16, bold=True, color=RED)
    add_text_box(slide, Inches(1.1), y + Inches(0.55), Inches(5), Inches(0.5),
                 desc, font_size=12, color=DARK)

add_rounded_card(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5.2),
                 RGBColor(0xF0, 0xFF, 0xF5), GREEN)
add_text_box(slide, Inches(7.4), Inches(1.8), Inches(4.8), Inches(0.6),
             "🎯  The Gap We Fill", font_size=22, bold=True, color=GREEN)
add_accent_line(slide, Inches(7.4), Inches(2.4), Inches(1.5))

gap_points = [
    "Zero AI-powered PCOS screening tools for Indian women",
    "100M+ women need early detection — today",
    "750M+ smartphone users in India",
    "Government actively pushing digital health",
    "Young women increasingly health-aware",
    "Preventive care saves lives and suffering",
]
for j, point in enumerate(gap_points):
    add_text_box(slide, Inches(7.4), Inches(2.7 + j * 0.55), Inches(4.8), Inches(0.5),
                 "→  " + point, font_size=14, color=DARK)


# ─── SLIDE 4: Our Solution ─────────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Introducing HormoCare AI", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
             "A comprehensive AI-powered PCOS early-detection and care platform",
             font_size=16, color=GRAY)

pillars = [
    ("📊  Smart Tracking", PINK,
     ["Menstrual cycle patterns & irregularity",
      "Weight changes & BMI trends",
      "Acne severity (1–5 visual scale)",
      "Hair fall & excess body hair",
      "Mood swings & sleep quality",
      "Family history integration"]),
    ("🧠  AI Risk Engine", TEAL,
     ["15+ symptom parameter analysis",
      "Weighted scoring algorithm",
      "Cycle: 40% • Weight: 20%",
      "Hormonal: 30% • Lifestyle: 10%",
      "Colour-coded risk (Green / Yellow / Red)",
      "Trend detection across weeks"]),
    ("💡  Personalised Care", GREEN,
     ["PCOS-friendly Indian diet plans",
      "Customised exercise routines",
      "Supplement guidance (Myo-inositol, Vit D)",
      "Doctor consultation timing alerts",
      "Mental health & stress support",
      "Educational health content"]),
]

for i, (title, accent, points) in enumerate(pillars):
    x = Inches(0.8 + i * 4.1)
    add_rounded_card(slide, x, Inches(2.1), Inches(3.8), Inches(4.8), LIGHT_BG)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, x, Inches(2.1), Inches(3.8), Inches(0.08)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = accent
    shape.line.fill.background()
    add_text_box(slide, x + Inches(0.3), Inches(2.3), Inches(3.2), Inches(0.5),
                 title, font_size=18, bold=True, color=accent)
    for j, pt in enumerate(points):
        add_text_box(slide, x + Inches(0.3), Inches(3.0 + j * 0.45),
                     Inches(3.2), Inches(0.4),
                     "▸  " + pt, font_size=12, color=DARK)


# ─── SLIDE 5: How the AI Works ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "How the AI Works", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
             "A rule-based intelligent system that analyses patterns no human eye can catch early enough.",
             font_size=16, color=GRAY)

flow_steps = [
    ("📱", "Data\nCollection", "User logs symptoms\ndaily in < 1 min"),
    ("⚙️", "Pattern\nDetection", "Algorithm spots\nirregularity trends"),
    ("🧠", "Risk\nScoring", "Weighted analysis\nacross 15+ factors"),
    ("🎯", "Risk\nLevel", "Green / Yellow / Red\nwith % score"),
    ("💡", "Action\nPlan", "Personalised next\nsteps & care"),
]
for i, (icon, title, desc) in enumerate(flow_steps):
    x = Inches(0.6 + i * 2.55)
    add_rounded_card(slide, x, Inches(2.1), Inches(2.2), Inches(2.4), LIGHT_BG, PINK)
    add_text_box(slide, x, Inches(2.2), Inches(2.2), Inches(0.6),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(1.8), Inches(0.6),
                 title, font_size=14, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.4), Inches(1.8), Inches(0.7),
                 desc, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)
    if i < len(flow_steps) - 1:
        add_text_box(slide, x + Inches(2.2), Inches(2.8), Inches(0.35), Inches(0.5),
                     "→", font_size=24, bold=True, color=PINK,
                     alignment=PP_ALIGN.CENTER)

add_rounded_card(slide, Inches(0.8), Inches(4.8), Inches(11.7), Inches(2.2),
                 RGBColor(0xF8, 0xF0, 0xFF))
add_text_box(slide, Inches(1.2), Inches(4.95), Inches(5), Inches(0.5),
             "Risk Algorithm Formula", font_size=18, bold=True, color=DARK_PINK)

formula_lines = [
    "PCOS Risk Score  =",
    "   (Cycle Irregularity Analysis  ×  40%)",
    " + (Weight & BMI Trend Analysis  ×  20%)",
    " + (Hormonal Symptoms Score        ×  30%)",
    " + (Lifestyle & Family Risk           ×  10%)",
]
for k, line in enumerate(formula_lines):
    add_text_box(slide, Inches(1.2), Inches(5.5 + k * 0.32), Inches(5.5), Inches(0.3),
                 line, font_size=13, color=DARK, font_name="Consolas")

examples = [
    ("HIGH RISK", "Cycle > 35 days  +  acne  +  weight gain  +  hair fall", RED),
    ("MODERATE", "Irregular cycle  +  hair fall  or  acne alone", ORANGE),
    ("LOW RISK", "Regular cycle  +  no hormonal symptoms", GREEN),
]
for k, (level, rule, clr) in enumerate(examples):
    y = Inches(5.05 + k * 0.65)
    add_rounded_card(slide, Inches(7.2), y, Inches(5), Inches(0.55), WHITE, clr)
    add_text_box(slide, Inches(7.4), y + Inches(0.05), Inches(1.2), Inches(0.4),
                 level, font_size=11, bold=True, color=clr)
    add_text_box(slide, Inches(8.7), y + Inches(0.05), Inches(3.3), Inches(0.4),
                 rule, font_size=11, color=DARK)


# ─── SLIDE 6: App Flow & Demo ──────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "User Journey — 4 Simple Steps", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

steps = [
    ("Step 1", "Profile Setup", PINK,
     ["Enter age, weight, height",
      "Family history of PCOS",
      "Current medications",
      "Lifestyle basics"]),
    ("Step 2", "Daily Symptom Log", TEAL,
     ["Period date & flow",
      "Acne level (1–5 scale)",
      "Hair fall (Low / Med / High)",
      "Mood & sleep hours"]),
    ("Step 3", "AI Risk Report", ORANGE,
     ["Risk score 0–100%",
      "Green / Yellow / Red indicator",
      "Key risk factors listed",
      "Trend graph over time"]),
    ("Step 4", "Care Plan", GREEN,
     ["Personalised diet chart",
      "Exercise recommendations",
      "Supplement reminders",
      "Doctor consultation alert"]),
]

for i, (step, title, accent, points) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    add_rounded_card(slide, x, Inches(1.6), Inches(2.9), Inches(5.2), LIGHT_BG)
    add_circle(slide, x + Inches(1.05), Inches(1.8), Inches(0.8), accent)
    add_text_box(slide, x + Inches(1.05), Inches(1.85), Inches(0.8), Inches(0.7),
                 str(i + 1), font_size=28, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(2.8), Inches(2.5), Inches(0.4),
                 step, font_size=12, bold=True, color=accent,
                 alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.15), Inches(2.5), Inches(0.5),
                 title, font_size=17, bold=True, color=DARK,
                 alignment=PP_ALIGN.CENTER)
    for j, pt in enumerate(points):
        add_text_box(slide, x + Inches(0.3), Inches(3.8 + j * 0.5),
                     Inches(2.3), Inches(0.4),
                     "✓  " + pt, font_size=12, color=DARK)
    if i < 3:
        add_text_box(slide, x + Inches(2.9), Inches(2.9), Inches(0.3), Inches(0.5),
                     "→", font_size=28, bold=True, color=PINK,
                     alignment=PP_ALIGN.CENTER)


# ─── SLIDE 7: Real-World Impact ────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8),
             "The Change We're Building", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
             "Impact that matters — measured in lives changed, not just numbers.",
             font_size=16, color=GRAY)

impacts = [
    ("🩺", "Early Detection Saves Lives",
     "Catching PCOS 2-3 years earlier means preventing irreversible damage — infertility, chronic diabetes, cardiovascular disease. One notification can change a woman's entire future.",
     RED),
    ("🧠", "Breaking the Silence",
     "Millions of women normalise irregular periods and suffer alone. HormoCare creates awareness at scale — teaching women that these symptoms are not 'normal' and help exists.",
     PINK),
    ("🌾", "Culturally Relevant Healthcare",
     "Built for Indian women — understanding local diets, family dynamics, and social barriers that prevent women from seeking help. Healthcare that speaks their language.",
     TEAL),
    ("👩‍⚕️", "Bridging Doctor Access Gaps",
     "In rural and semi-urban India, gynaecologist access is limited. AI-powered pre-screening helps women know when professional help is truly urgent, reducing unnecessary delays.",
     GREEN),
    ("💪", "Empowering Self-Advocacy",
     "Armed with data and risk reports, women can walk into a doctor's office with evidence — not dismissed with 'it's just stress'. Data gives them a voice.",
     ORANGE),
    ("🤝", "Reducing Family Burden",
     "Untreated PCOS leads to fertility issues, mental health struggles, and long-term disease. Early intervention protects not just the woman, but her entire family's wellbeing.",
     BLUE),
]

for i, (icon, title, desc, accent) in enumerate(impacts):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(2.0 + row * 1.7)
    add_rounded_card(slide, x, y, Inches(5.9), Inches(1.5), WHITE, accent)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.15), Inches(0.5), Inches(0.5),
                 icon, font_size=22)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.12), Inches(4.9), Inches(0.35),
                 title, font_size=15, bold=True, color=accent)
    add_text_box(slide, x + Inches(0.7), y + Inches(0.5), Inches(4.9), Inches(0.9),
                 desc, font_size=11, color=DARK, line_spacing=1.4)


# ─── SLIDE 8: Innovation Features ──────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Innovation That Sets Us Apart", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

innovations = [
    ("📊", "Cycle Irregularity Graph",
     "Visual trend showing cycle patterns over months — making the invisible visible for both users and doctors.",
     PINK),
    ("🤖", "Health Chatbot",
     "24/7 conversational AI answering hormonal health questions, reducing anxiety and spreading awareness.",
     TEAL),
    ("🧠", "Mood–Cycle Correlation",
     "Maps mood patterns against menstrual cycle phases — revealing hormonal connections women never knew existed.",
     ORANGE),
    ("🥗", "PCOS Diet Planner",
     "Indian-cuisine focused meal plans: low GI rotis, millets, anti-inflammatory sabzis — not Western kale smoothies.",
     GREEN),
    ("🏥", "Gynecologist Finder",
     "Locate PCOS-specialist doctors nearby with ratings, availability, and direct appointment booking.",
     BLUE),
    ("📱", "Shareable Reports",
     "Generate PDF health reports to share with doctors — no more explaining months of symptoms in a 5-min consultation.",
     DARK_PINK),
]

for i, (icon, title, desc, accent) in enumerate(innovations):
    col = i % 3
    row = i // 3
    x = Inches(0.6 + col * 4.1)
    y = Inches(1.6 + row * 2.8)
    add_rounded_card(slide, x, y, Inches(3.8), Inches(2.5), LIGHT_BG, accent)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(3.2), Inches(0.5),
                 icon + "  " + title, font_size=16, bold=True, color=accent)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.2), Inches(1.5),
                 desc, font_size=12, color=DARK, line_spacing=1.5)


# ─── SLIDE 9: Technical Architecture ───────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Technical Architecture", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))

layers = [
    ("Presentation Layer", PINK,
     "HTML5 / CSS3 / JavaScript\nResponsive Mobile-First UI\nProgressive Web App (PWA)"),
    ("Application Layer", TEAL,
     "Rule-Based AI Risk Engine\nSymptom Pattern Analyser\nRecommendation Generator"),
    ("Data Layer", ORANGE,
     "Local Storage (Privacy-First)\nEncrypted User Profiles\nSymptom History & Trends"),
    ("Integration Layer", GREEN,
     "Doctor Directory API\nDiet & Exercise Database\nNotification System"),
]

for i, (title, accent, desc) in enumerate(layers):
    y = Inches(1.6 + i * 1.35)
    add_rounded_card(slide, Inches(0.8), y, Inches(5.5), Inches(1.15), LIGHT_BG, accent)
    add_text_box(slide, Inches(1.1), y + Inches(0.15), Inches(2), Inches(0.4),
                 title, font_size=14, bold=True, color=accent)
    add_text_box(slide, Inches(3.2), y + Inches(0.15), Inches(2.8), Inches(0.85),
                 desc, font_size=11, color=DARK)

add_rounded_card(slide, Inches(7), Inches(1.6), Inches(5.5), Inches(5.5),
                 RGBColor(0xF8, 0xF0, 0xFF))
add_text_box(slide, Inches(7.4), Inches(1.8), Inches(4.8), Inches(0.5),
             "🛡️  Privacy & Security First", font_size=18, bold=True, color=DARK_PINK)
add_accent_line(slide, Inches(7.4), Inches(2.35), Inches(1.5))

privacy_points = [
    "All data stored locally on user's device",
    "Zero data transmission without explicit consent",
    "No third-party tracking or analytics",
    "User can export or delete all data anytime",
    "HIPAA-aligned data handling principles",
    "End-to-end encryption for cloud sync (future)",
]
for j, pt in enumerate(privacy_points):
    add_text_box(slide, Inches(7.4), Inches(2.6 + j * 0.55), Inches(4.8), Inches(0.5),
                 "🔒  " + pt, font_size=13, color=DARK)

add_text_box(slide, Inches(7.4), Inches(5.95), Inches(4.8), Inches(0.7),
             "Scalable to: Flutter / React Native mobile apps\n"
             "Future: Firebase backend + ML model upgrade",
             font_size=12, color=GRAY)


# ─── SLIDE 10: Future Vision ───────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.8),
             "Where This Goes Next", font_size=38, bold=True, color=DARK)
add_accent_line(slide, Inches(0.8), Inches(1.15), Inches(2))
add_text_box(slide, Inches(0.8), Inches(1.35), Inches(11), Inches(0.5),
             "From a prototype to a platform that transforms women's healthcare across India and beyond.",
             font_size=16, color=GRAY)

phases = [
    ("Phase 1", "Foundation", "3–6 Months", PINK,
     ["Core app with AI risk assessment",
      "Symptom tracking & recommendations",
      "Community testing & feedback",
      "Validation with gynecologists"]),
    ("Phase 2", "Healthcare Integration", "6–12 Months", TEAL,
     ["Doctor appointment booking",
      "Lab test partnerships",
      "Pharmacy integration",
      "Regional language support"]),
    ("Phase 3", "Platform Scale", "12–24 Months", GREEN,
     ["Telemedicine consultations",
      "Machine learning model upgrade",
      "Wearable device integration",
      "Government health scheme tie-ups"]),
]

for i, (phase, title, timeline, accent, points) in enumerate(phases):
    x = Inches(0.6 + i * 4.15)
    add_rounded_card(slide, x, Inches(2.0), Inches(3.85), Inches(4.4), LIGHT_BG, accent)
    add_text_box(slide, x + Inches(0.3), Inches(2.15), Inches(3.3), Inches(0.4),
                 phase + " — " + timeline, font_size=12, bold=True, color=accent)
    add_text_box(slide, x + Inches(0.3), Inches(2.55), Inches(3.3), Inches(0.5),
                 title, font_size=20, bold=True, color=DARK)
    for j, pt in enumerate(points):
        add_text_box(slide, x + Inches(0.3), Inches(3.3 + j * 0.55),
                     Inches(3.3), Inches(0.45),
                     "→  " + pt, font_size=13, color=DARK)

add_rounded_card(slide, Inches(0.6), Inches(6.6), Inches(12.1), Inches(0.7),
                 RGBColor(0xF0, 0xF8, 0xFF), BLUE)
add_text_box(slide, Inches(1), Inches(6.65), Inches(11.5), Inches(0.55),
             "🌏  Long-term: Expand to Southeast Asia & Middle East — regions with similar PCOS prevalence and smartphone growth",
             font_size=14, color=DARK, alignment=PP_ALIGN.CENTER)


# ─── SLIDE 11: Thank You / CTA ─────────────────────────────────────────
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, 0xFF, 0xFF, 0xFF)
add_gradient_bar(slide, top=0)
add_gradient_bar(slide, top=H - Inches(0.08))

add_circle(slide, Inches(10), Inches(0.5), Inches(5), RGBColor(0xFF, 0xF0, 0xF5))
add_circle(slide, Inches(-1.5), Inches(4), Inches(4), RGBColor(0xFF, 0xF0, 0xF5))

add_text_box(slide, Inches(1.5), Inches(1.0), Inches(10), Inches(0.8),
             "🌸", font_size=50, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1),
             "HormoCare AI", font_size=48, bold=True, color=PINK,
             alignment=PP_ALIGN.CENTER)

add_accent_line(slide, Inches(5.5), Inches(2.85), Inches(2.3))

add_rich_text_box(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(1.5), [
    {"text": '"Every woman deserves to know her body.',
     "size": 22, "color": DARK, "align": PP_ALIGN.CENTER, "space_after": 4},
    {"text": 'Every symptom deserves to be heard.',
     "size": 22, "color": DARK, "align": PP_ALIGN.CENTER, "space_after": 4},
    {"text": 'Every diagnosis deserves to come early."',
     "size": 22, "color": DARK, "align": PP_ALIGN.CENTER, "space_after": 16},
    {"text": "Let's build preventive healthcare — together.",
     "size": 18, "bold": True, "color": PINK, "align": PP_ALIGN.CENTER,
     "space_after": 0},
])

add_text_box(slide, Inches(3), Inches(5.5), Inches(7), Inches(0.5),
             "Thank you for your time.",
             font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(3), Inches(6.2), Inches(7), Inches(0.5),
             "Questions & Discussion",
             font_size=20, bold=True, color=DARK, alignment=PP_ALIGN.CENTER)


# ─── SAVE ──────────────────────────────────────────────────────────────
output_path = "/Users/mmt12368/Desktop/vtpaurelia/hormocare-ai/HormoCare_AI_Pitch.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
