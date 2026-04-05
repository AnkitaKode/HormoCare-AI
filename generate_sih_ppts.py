from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──────────────────────────────────────────────────────
PINK      = RGBColor(0xFF, 0x6B, 0x9D)
DARK_PINK = RGBColor(0xC4, 0x45, 0x69)
DARK      = RGBColor(0x2C, 0x3E, 0x50)
GRAY      = RGBColor(0x7F, 0x8C, 0x8D)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF8, 0xF9, 0xFA)
GREEN     = RGBColor(0x2E, 0xCC, 0x71)
RED       = RGBColor(0xE7, 0x4C, 0x3C)
ORANGE    = RGBColor(0xF3, 0x9C, 0x12)
TEAL      = RGBColor(0x4E, 0xCD, 0xC4)
BLUE      = RGBColor(0x45, 0xB7, 0xD1)
NAVY      = RGBColor(0x1A, 0x1A, 0x2E)
DEEP_BLUE = RGBColor(0x16, 0x21, 0x3E)
MID_BLUE  = RGBColor(0x0F, 0x3D, 0x6E)
BLACK     = RGBColor(0x00, 0x00, 0x00)
LGRAY     = RGBColor(0xBD, 0xC3, 0xC7)
TABLE_HDR = RGBColor(0x2C, 0x3E, 0x50)
TABLE_ALT = RGBColor(0xEC, 0xF0, 0xF1)

# ── Helpers ────────────────────────────────────────────────────────────
def set_bg(slide, r, g, b):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def top_bar(slide, W):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = PINK; s.line.fill.background()


def bottom_bar(slide, W, H):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.5), W, Inches(0.5))
    s.fill.solid(); s.fill.fore_color.rgb = DEEP_BLUE; s.line.fill.background()


def footer_text(slide, W, H, text="Team Sparsh  |  HormoCare AI  |  SIH 2025"):
    txb = slide.shapes.add_textbox(Inches(0.3), H - Inches(0.42), W - Inches(0.6), Inches(0.35))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER


def accent_line(slide, left, top, width):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Inches(0.05))
    s.fill.solid(); s.fill.fore_color.rgb = PINK; s.line.fill.background()


def card(slide, left, top, width, height, fill=WHITE, border=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if border:
        s.line.color.rgb = border; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s


def txt(slide, left, top, width, height, text, size=14, bold=False,
        color=DARK, align=PP_ALIGN.LEFT, font="Arial", spacing=1.2):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    p.font.name = font; p.alignment = align
    p.space_after = Pt(0)
    return txb


def rich(slide, left, top, width, height, lines, font="Arial"):
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf = txb.text_frame; tf.word_wrap = True
    for i, ld in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = ld.get("t", "")
        p.font.size = Pt(ld.get("s", 14))
        p.font.bold = ld.get("b", False)
        p.font.color.rgb = ld.get("c", DARK)
        p.font.name = font
        p.alignment = ld.get("a", PP_ALIGN.LEFT)
        p.space_after = Pt(ld.get("sa", 4))
        p.space_before = Pt(ld.get("sb", 0))
    return txb


def slide_title(slide, title, subtitle=None, W=None):
    txt(slide, Inches(0.8), Inches(0.3), Inches(9), Inches(0.7),
        title, size=30, bold=True, color=DARK)
    accent_line(slide, Inches(0.8), Inches(0.95), Inches(1.8))
    if subtitle:
        txt(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
            subtitle, size=13, color=GRAY)


# ======================================================================
#  PART 1 : IDEATION ROUND PPT (6 SLIDES)
# ======================================================================

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ─── SLIDE 1: Team Details & Problem Statement ────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)

# Header bar
hdr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.06), W, Inches(1.05))
hdr.fill.solid(); hdr.fill.fore_color.rgb = DEEP_BLUE; hdr.line.fill.background()

txt(sl, Inches(0.8), Inches(0.15), Inches(11), Inches(0.5),
    "SIH 2025 — Ideation Round Submission", size=14, color=LGRAY, align=PP_ALIGN.CENTER)
txt(sl, Inches(0.8), Inches(0.5), Inches(11), Inches(0.55),
    "Team Details & Problem Statement", size=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left card — Problem Statement
card(sl, Inches(0.6), Inches(1.4), Inches(6), Inches(2.4), LIGHT_BG, PINK)
rich(sl, Inches(0.9), Inches(1.5), Inches(5.4), Inches(2.2), [
    {"t": "Problem Statement", "s": 18, "b": True, "c": PINK, "sa": 10},
    {"t": "Problem Statement ID:  _____________  (to be assigned)", "s": 13, "c": DARK, "sa": 6},
    {"t": "Problem Statement Title:", "s": 13, "b": True, "c": DARK, "sa": 4},
    {"t": "AI-Powered Early PCOS Detection & Personalised Care", "s": 14, "b": True, "c": DARK_PINK, "sa": 4},
    {"t": "Platform for Indian Women — \"HormoCare AI\"", "s": 14, "b": True, "c": DARK_PINK, "sa": 8},
    {"t": "Category: Software  |  Domain: MedTech / Women's Health / AI", "s": 11, "c": GRAY, "sa": 0},
])

# Right card — Team Details
card(sl, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.4), LIGHT_BG, TEAL)
rich(sl, Inches(7.2), Inches(1.5), Inches(5.2), Inches(2.2), [
    {"t": "Team Details", "s": 18, "b": True, "c": TEAL, "sa": 10},
    {"t": "Team Name:  Sparsh", "s": 13, "b": True, "c": DARK, "sa": 4},
    {"t": "Team Leader:  Anshika Madhu", "s": 13, "b": True, "c": DARK, "sa": 4},
    {"t": "Institution:  MITMAAI", "s": 13, "c": DARK, "sa": 0},
])

# Team Members Table
members = [
    ("S.No", "Name", "Reg. No.", "Role"),
    ("1", "Anshika Madhu", "25P108", "Team Leader"),
    ("2", "Nilakshi", "25P84", "Member"),
    ("3", "Anchal Kumari", "25P70", "Member"),
    ("4", "Mayank Kumar", "25P27", "Member"),
    ("5", "Arpan Gandhi", "25P14", "Member"),
    ("6", "Saurabh Kumar", "25P49", "Member"),
]

table_top = Inches(4.05)
rows = len(members)
cols = 4
tbl_shape = sl.shapes.add_table(rows, cols, Inches(0.6), table_top, Inches(8.5), Inches(2.5))
tbl = tbl_shape.table

col_widths = [Inches(0.7), Inches(3.0), Inches(2.0), Inches(2.8)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = members[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            else:
                p.font.color.rgb = DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HDR
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE

# Mentor info
card(sl, Inches(9.5), Inches(4.05), Inches(3.2), Inches(2.5), LIGHT_BG, ORANGE)
rich(sl, Inches(9.7), Inches(4.15), Inches(2.8), Inches(2.3), [
    {"t": "Mentor Details", "s": 14, "b": True, "c": ORANGE, "sa": 8},
    {"t": "Name: _______________", "s": 12, "c": DARK, "sa": 4},
    {"t": "Branch: _______________", "s": 12, "c": DARK, "sa": 4},
    {"t": "Designation: _______________", "s": 12, "c": DARK, "sa": 4},
    {"t": "Domain: MedTech / AI / CS", "s": 12, "c": DARK, "sa": 4},
    {"t": "(Same domain as project)", "s": 10, "c": GRAY, "sa": 0},
])


# ─── SLIDE 2: Proposed Solution ──────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)
slide_title(sl, "Proposed Solution")

# Core Idea bullets
card(sl, Inches(0.6), Inches(1.5), Inches(7.2), Inches(3.5), LIGHT_BG, PINK)
txt(sl, Inches(0.9), Inches(1.6), Inches(6.6), Inches(0.4),
    "Core Idea", size=16, bold=True, color=PINK)

bullets = [
    "An AI-powered web platform that enables early PCOS detection by tracking and analysing 15+ health parameters — menstrual cycle, weight, acne, hair fall, mood, and sleep patterns.",
    "Uses a weighted rule-based AI algorithm to generate a personalised PCOS Risk Score (0–100%) with colour-coded severity (Green / Yellow / Red) and trend graphs.",
    "Delivers actionable care plans: PCOS-friendly Indian diet recommendations, exercise routines, supplement guidance (Myo-inositol, Vitamin D), and doctor consultation timing alerts.",
    "Designed specifically for Indian women — considering local dietary habits, genetic predisposition, cultural barriers, and healthcare access gaps in rural/semi-urban areas.",
]
for i, b in enumerate(bullets):
    rich(sl, Inches(0.9), Inches(2.15 + i * 0.7), Inches(6.6), Inches(0.65), [
        {"t": f"▸  {b}", "s": 12, "c": DARK, "sa": 2},
    ])

# USP card
card(sl, Inches(8.1), Inches(1.5), Inches(4.6), Inches(3.5), RGBColor(0xFF, 0xF0, 0xF5), DARK_PINK)
txt(sl, Inches(8.4), Inches(1.6), Inches(4), Inches(0.4),
    "Unique Selling Proposition (USP)", size=16, bold=True, color=DARK_PINK)

usps = [
    ("🎯", "First AI-Based PCOS Screener", "for Indian women specifically"),
    ("🧠", "15+ Parameter AI Analysis", "not just a period tracker"),
    ("🇮🇳", "India-Specific Design", "local diets, genetics, culture"),
    ("🔒", "Privacy-First Architecture", "all data stays on device"),
    ("⚡", "Works Offline / Low Data", "accessible in rural areas"),
    ("📊", "Shareable Doctor Reports", "evidence-based consultations"),
]
for i, (icon, title, sub) in enumerate(usps):
    y = Inches(2.15 + i * 0.45)
    rich(sl, Inches(8.4), y, Inches(4), Inches(0.42), [
        {"t": f"{icon}  {title}", "s": 12, "b": True, "c": DARK, "sa": 0},
        {"t": f"      {sub}", "s": 10, "c": GRAY, "sa": 0},
    ])

# Bottom — impact statement
card(sl, Inches(0.6), Inches(5.2), Inches(12.1), Inches(1.5), RGBColor(0xF0, 0xFF, 0xF5), GREEN)
txt(sl, Inches(0.9), Inches(5.3), Inches(11.5), Inches(0.4),
    "Why This Matters", size=16, bold=True, color=GREEN)
rich(sl, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.8), [
    {"t": "1 in 5 Indian women has PCOS. 70% go undiagnosed for 2–3 years. This leads to infertility, diabetes, obesity, and depression.", "s": 12, "c": DARK, "sa": 4},
    {"t": "HormoCare AI bridges this gap by putting early detection directly in women's hands — no expensive tests, no waiting months for a specialist.", "s": 12, "b": True, "c": DARK, "sa": 0},
])


# ─── SLIDE 3: Technical Stack ────────────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)
slide_title(sl, "Technical Stack")

stacks = [
    ("Frontend / UI", "🖥️", PINK, [
        "HTML5 — Semantic, accessible markup",
        "CSS3 — Grid, Flexbox, responsive design",
        "JavaScript (ES6+) — Interactive UI logic",
        "Progressive Web App (PWA) — Installable on mobile",
        "Font Awesome — Icon library",
        "Google Fonts (Inter) — Typography",
    ]),
    ("AI / Logic Engine", "🧠", TEAL, [
        "Rule-Based AI Algorithm — Weighted risk scoring",
        "Pattern Detection — Cycle irregularity analysis",
        "Multi-Parameter Analysis — 15+ symptom factors",
        "Trend Computation — Weekly/monthly patterns",
        "Risk Classification — Green / Yellow / Red",
        "Recommendation Engine — Personalised output",
    ]),
    ("Data / Storage", "💾", ORANGE, [
        "localStorage — Client-side data persistence",
        "JSON — Structured symptom data format",
        "Data Export — PDF / CSV report generation",
        "Encryption — AES for sensitive health data",
        "Zero Cloud Dependency — Works fully offline",
        "HIPAA-Aligned Principles — Privacy first",
    ]),
    ("Future Scale Stack", "🚀", GREEN, [
        "React Native / Flutter — Mobile app",
        "Firebase — Auth, Firestore, Cloud Functions",
        "TensorFlow.js — ML model upgrade",
        "Node.js + Express — Backend API",
        "MongoDB — Scalable health data storage",
        "Twilio / FCM — Notifications & reminders",
    ]),
]

for i, (title, icon, accent, items) in enumerate(stacks):
    col = i % 2
    row = i // 2
    x = Inches(0.6 + col * 6.2)
    y = Inches(1.4 + row * 2.85)
    card(sl, x, y, Inches(5.9), Inches(2.65), LIGHT_BG, accent)
    # Header stripe
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(5.9), Inches(0.45))
    s.fill.solid(); s.fill.fore_color.rgb = accent; s.line.fill.background()
    txt(sl, x + Inches(0.2), y + Inches(0.05), Inches(5.4), Inches(0.35),
        f"{icon}  {title}", size=14, bold=True, color=WHITE)
    for j, item in enumerate(items):
        txt(sl, x + Inches(0.3), y + Inches(0.55 + j * 0.33), Inches(5.3), Inches(0.3),
            f"▸  {item}", size=11, color=DARK)


# ─── SLIDE 4: Architecture Diagram ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)
slide_title(sl, "System Architecture & Data Flow")

# 3-Layer Architecture
layers = [
    ("PRESENTATION LAYER (User Interface)", PINK,
     "📱 Responsive Web App (HTML/CSS/JS)  →  PWA Installable  →  Mobile-First Design\n"
     "User interactions: Profile Setup  |  Symptom Logging  |  Risk Dashboard  |  Care Plans"),
    ("APPLICATION LAYER (Business Logic & AI)", TEAL,
     "🧠 AI Risk Engine: Weighted Algorithm (Cycle 40% + Weight 20% + Hormonal 30% + Lifestyle 10%)\n"
     "Pattern Detector  →  Risk Classifier  →  Recommendation Generator  →  Alert System"),
    ("DATA LAYER (Storage & Security)", ORANGE,
     "💾 localStorage (JSON)  →  AES Encryption  →  No Cloud Transfer Without Consent\n"
     "User Profile  |  Symptom History  |  Risk Assessments  |  Care Plan Records"),
]

for i, (title, accent, desc) in enumerate(layers):
    y = Inches(1.4 + i * 1.3)
    card(sl, Inches(0.6), y, Inches(12.1), Inches(1.15), WHITE, accent)
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), y, Inches(0.12), Inches(1.15))
    s.fill.solid(); s.fill.fore_color.rgb = accent; s.line.fill.background()
    txt(sl, Inches(1.0), y + Inches(0.08), Inches(11.4), Inches(0.3),
        title, size=13, bold=True, color=accent)
    txt(sl, Inches(1.0), y + Inches(0.4), Inches(11.4), Inches(0.65),
        desc, size=11, color=DARK)

# Arrows between layers
for i in range(2):
    y = Inches(2.55 + i * 1.3)
    txt(sl, Inches(5.8), y, Inches(1.5), Inches(0.35),
        "⬇  ⬆", size=16, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Data Flow Diagram
card(sl, Inches(0.6), Inches(5.15), Inches(12.1), Inches(1.6), RGBColor(0xF5, 0xF0, 0xFF))
txt(sl, Inches(0.9), Inches(5.2), Inches(5), Inches(0.3),
    "Data Flow Diagram (DFD)", size=14, bold=True, color=DARK_PINK)

flow_items = [
    ("👤\nUser", PINK),
    ("→", GRAY),
    ("📱\nApp UI\n(Input)", PINK),
    ("→", GRAY),
    ("⚙️\nValidation\n& Parsing", TEAL),
    ("→", GRAY),
    ("🧠\nAI Risk\nEngine", TEAL),
    ("→", GRAY),
    ("📊\nRisk Score\n& Report", ORANGE),
    ("→", GRAY),
    ("💡\nCare Plan\nOutput", GREEN),
    ("→", GRAY),
    ("👤\nUser\n(Result)", PINK),
]
for i, (label, clr) in enumerate(flow_items):
    x = Inches(0.7 + i * 0.95)
    if label == "→":
        txt(sl, x, Inches(5.75), Inches(0.6), Inches(0.4),
            "→", size=20, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
    else:
        card(sl, x, Inches(5.55), Inches(0.85), Inches(1.0), LIGHT_BG, clr)
        txt(sl, x, Inches(5.6), Inches(0.85), Inches(0.9),
            label, size=9, bold=True, color=clr, align=PP_ALIGN.CENTER)


# ─── SLIDE 5: Feasibility & Impact ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)
slide_title(sl, "Feasibility & Potential Impact")

# End Users
card(sl, Inches(0.6), Inches(1.4), Inches(6), Inches(2.3), LIGHT_BG, PINK)
txt(sl, Inches(0.9), Inches(1.5), Inches(5.4), Inches(0.35),
    "🎯  Who Are the End Users?", size=16, bold=True, color=PINK)

users = [
    "Women aged 15–40 experiencing menstrual irregularities or PCOS symptoms",
    "Rural/semi-urban women with limited access to gynaecologists",
    "Health-conscious young women wanting preventive screening",
    "Gynaecologists seeking pre-screening data from patients",
    "Corporate wellness programs for female employee health",
]
for i, u in enumerate(users):
    txt(sl, Inches(0.9), Inches(2.0 + i * 0.35), Inches(5.4), Inches(0.3),
        f"▸  {u}", size=11, color=DARK)

# Scale across India
card(sl, Inches(6.9), Inches(1.4), Inches(5.8), Inches(2.3), LIGHT_BG, GREEN)
txt(sl, Inches(7.2), Inches(1.5), Inches(5.2), Inches(0.35),
    "🇮🇳  How Will This Scale Across India?", size=16, bold=True, color=GREEN)

scale_pts = [
    "PWA works on any smartphone — no app store dependency",
    "Offline-first: functions without internet in rural areas",
    "Multilingual support roadmap (Hindi, Tamil, Bengali, etc.)",
    "Partnership with ASHA workers for community outreach",
    "Integration with Ayushman Bharat Digital Mission (ABDM)",
]
for i, s in enumerate(scale_pts):
    txt(sl, Inches(7.2), Inches(2.0 + i * 0.35), Inches(5.2), Inches(0.3),
        f"▸  {s}", size=11, color=DARK)

# Impact metrics
card(sl, Inches(0.6), Inches(3.95), Inches(12.1), Inches(2.7), WHITE, TEAL)
txt(sl, Inches(0.9), Inches(4.05), Inches(5), Inches(0.35),
    "📊  Potential Impact Metrics", size=16, bold=True, color=TEAL)

impacts = [
    ("🩺", "Early Detection", "Reduce diagnosis delay from 2–3 years to under 6 months through proactive screening alerts", RED),
    ("🧠", "Awareness at Scale", "Educate millions that irregular periods + acne + hair fall are not 'normal' but PCOS warning signs", PINK),
    ("💪", "Women's Empowerment", "Arm women with data-backed health reports so doctors take their symptoms seriously", ORANGE),
    ("🏥", "Healthcare Access", "Bridge the gynaecologist gap in rural India through AI-powered pre-screening", GREEN),
    ("🤝", "Family Impact", "Prevent downstream effects — infertility, chronic diabetes, depression — improving entire family wellbeing", BLUE),
]
for i, (icon, title, desc, clr) in enumerate(impacts):
    y = Inches(4.5 + i * 0.42)
    txt(sl, Inches(0.9), y, Inches(11.5), Inches(0.38),
        f"{icon}  {title}:  {desc}", size=11, color=DARK)


# ─── SLIDE 6: Photos & Attestation ──────────────────────────────────
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar(sl, W); bottom_bar(sl, W, H)
footer_text(sl, W, H)
slide_title(sl, "Team Photos & Attestation")

# Photo placeholders
members_list = [
    ("Anshika Madhu", "25P108", "Team Leader"),
    ("Nilakshi", "25P84", "Member"),
    ("Anchal Kumari", "25P70", "Member"),
    ("Mayank Kumar", "25P27", "Member"),
    ("Arpan Gandhi", "25P14", "Member"),
    ("Saurabh Kumar", "25P49", "Member"),
]

for i, (name, reg, role) in enumerate(members_list):
    col = i % 6
    x = Inches(0.5 + col * 2.1)
    y = Inches(1.4)
    # Photo placeholder box
    card(sl, x, y, Inches(1.8), Inches(2.0), LIGHT_BG, LGRAY)
    txt(sl, x, y + Inches(0.6), Inches(1.8), Inches(0.5),
        "📷", size=28, align=PP_ALIGN.CENTER, color=LGRAY)
    txt(sl, x, y + Inches(1.15), Inches(1.8), Inches(0.3),
        "[Photo]", size=10, color=LGRAY, align=PP_ALIGN.CENTER)
    # Name below
    txt(sl, x, y + Inches(2.05), Inches(1.8), Inches(0.25),
        name, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, x, y + Inches(2.3), Inches(1.8), Inches(0.2),
        f"{reg}  |  {role}", size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Attestation
card(sl, Inches(0.5), Inches(4.1), Inches(12.3), Inches(2.6), RGBColor(0xFF, 0xFB, 0xF0), ORANGE)
txt(sl, Inches(0.8), Inches(4.2), Inches(11.7), Inches(0.35),
    "📜  Declaration & Attestation", size=16, bold=True, color=ORANGE)

attestation = (
    "We agree to abide by all rules and policies set forth by MITMAAI. "
    "The decision of MITMAAI selection team will be binding and cannot be challenged.\n\n"
    "We also agree that, before filing any patent applications or forming any startup "
    "based substantially on their winning project, we will first consult with MITMAAI/the Mentor "
    "regarding strategy, protection, and commercialization, and will consider in good faith "
    "offering MITMAAI/the Mentor an appropriate advisory or equity role as may be mutually agreed in writing."
)
txt(sl, Inches(0.8), Inches(4.6), Inches(11.7), Inches(1.6),
    attestation, size=11, color=DARK)

# Signature line
sig_names = ["Anshika Madhu", "Nilakshi", "Anchal Kumari", "Mayank Kumar", "Arpan Gandhi", "Saurabh Kumar"]
for i, name in enumerate(sig_names):
    x = Inches(0.6 + i * 2.1)
    y = Inches(6.4)
    txt(sl, x, y, Inches(1.8), Inches(0.15),
        "______________________", size=9, color=LGRAY, align=PP_ALIGN.CENTER)
    txt(sl, x, y + Inches(0.18), Inches(1.8), Inches(0.2),
        name, size=9, color=DARK, align=PP_ALIGN.CENTER)


# ── Save Ideation PPT ─────────────────────────────────────────────────
path1 = "/Users/mmt12368/Desktop/vtpaurelia/hormocare-ai/HormoCare_AI_SIH_Ideation.pptx"
prs.save(path1)
print(f"✓ Ideation PPT saved: {path1}  ({len(prs.slides)} slides)")


# ======================================================================
#  PART 2 : GRAND FINALE PPT (8 SLIDES)
# ======================================================================

prs2 = Presentation()
prs2.slide_width = Inches(13.333)
prs2.slide_height = Inches(7.5)
W = prs2.slide_width
H = prs2.slide_height

def top_bar2(sl):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    s.fill.solid(); s.fill.fore_color.rgb = PINK; s.line.fill.background()
def bottom_bar2(sl):
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(0.5), W, Inches(0.5))
    s.fill.solid(); s.fill.fore_color.rgb = DEEP_BLUE; s.line.fill.background()
def footer2(sl, t="Team Sparsh  |  HormoCare AI  |  Grand Finale"):
    txb = sl.shapes.add_textbox(Inches(0.3), H - Inches(0.42), W - Inches(0.6), Inches(0.35))
    tf = txb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = t
    p.font.size = Pt(10); p.font.color.rgb = WHITE; p.font.name = "Arial"
    p.alignment = PP_ALIGN.CENTER
def stitle(sl, t, sub=None):
    txt(sl, Inches(0.8), Inches(0.3), Inches(9), Inches(0.7),
        t, size=30, bold=True, color=DARK)
    accent_line(sl, Inches(0.8), Inches(0.95), Inches(1.8))
    if sub:
        txt(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.4),
            sub, size=13, color=GRAY)


# ─── F-SLIDE 1: Title + Team ─────────────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)

hdr = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.06), W, Inches(2.2))
hdr.fill.solid(); hdr.fill.fore_color.rgb = DEEP_BLUE; hdr.line.fill.background()

txt(sl, Inches(1), Inches(0.3), Inches(11), Inches(0.5),
    "SIH 2025  —  Grand Finale", size=14, color=LGRAY, align=PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(0.8), Inches(11), Inches(0.7),
    "🌸  HormoCare AI", size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(sl, Inches(1), Inches(1.5), Inches(11), Inches(0.5),
    "AI-Powered Early PCOS Detection & Personalised Care for Indian Women",
    size=18, color=LGRAY, align=PP_ALIGN.CENTER)

txt(sl, Inches(1), Inches(2.6), Inches(11), Inches(0.4),
    "Team Sparsh  |  Leader: Anshika Madhu",
    size=16, bold=True, color=DARK, align=PP_ALIGN.CENTER)

m_names = ["Anshika Madhu\n(25P108, Leader)", "Nilakshi\n(25P84)", "Anchal Kumari\n(25P70)",
           "Mayank Kumar\n(25P27)", "Arpan Gandhi\n(25P14)", "Saurabh Kumar\n(25P49)"]
for i, nm in enumerate(m_names):
    x = Inches(0.6 + i * 2.1)
    card(sl, x, Inches(3.2), Inches(1.9), Inches(1.3), LIGHT_BG, PINK if i == 0 else LGRAY)
    txt(sl, x, Inches(3.35), Inches(1.9), Inches(1.0),
        nm, size=11, bold=(i == 0), color=DARK, align=PP_ALIGN.CENTER)

# Quick stats
stats = [("20%", "Indian women have PCOS"), ("70%", "Go undiagnosed"), ("100M+", "Women affected"), ("2–3 yrs", "Diagnosis delay")]
for i, (num, lab) in enumerate(stats):
    x = Inches(0.8 + i * 3.1)
    card(sl, x, Inches(4.8), Inches(2.8), Inches(1.4), LIGHT_BG, PINK)
    txt(sl, x, Inches(4.9), Inches(2.8), Inches(0.6),
        num, size=30, bold=True, color=PINK, align=PP_ALIGN.CENTER)
    txt(sl, x, Inches(5.5), Inches(2.8), Inches(0.5),
        lab, size=12, color=DARK, align=PP_ALIGN.CENTER)


# ─── F-SLIDE 2: The Prototype ────────────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "The Prototype", "High-quality screenshots of our working platform")

screens = [
    ("🏠 Landing Page", "Hero section with PCOS statistics,\nproblem statement, feature showcase,\nand 'Start Assessment' CTA"),
    ("📝 Profile Setup", "User onboarding: age, weight,\nheight, family history input\nwith clean form design"),
    ("📊 Symptom Tracker", "Daily logging interface:\nmenstrual cycle, acne (1–5),\nhair fall, mood, sleep"),
    ("🧠 AI Risk Report", "Risk score (0–100%) with\ncolour indicator, risk factors,\nAI insights & trend analysis"),
    ("💡 Care Plan", "Personalised diet, exercise,\nsupplement reminders,\ndoctor consultation alerts"),
]

for i, (title, desc) in enumerate(screens):
    x = Inches(0.4 + i * 2.55)
    # Screenshot placeholder
    card(sl, x, Inches(1.5), Inches(2.3), Inches(3.2), LIGHT_BG, PINK)
    txt(sl, x, Inches(2.3), Inches(2.3), Inches(0.5),
        "📱", size=36, align=PP_ALIGN.CENTER, color=LGRAY)
    txt(sl, x, Inches(2.9), Inches(2.3), Inches(0.4),
        "[Screenshot]", size=10, color=LGRAY, align=PP_ALIGN.CENTER)
    # Label
    txt(sl, x, Inches(4.8), Inches(2.3), Inches(0.35),
        title, size=12, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(sl, x + Inches(0.1), Inches(5.2), Inches(2.1), Inches(1.0),
        desc, size=10, color=GRAY, align=PP_ALIGN.CENTER)


# ─── F-SLIDE 3: Detailed Workflow ────────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "Detailed Workflow", "Step-by-step: how HormoCare AI works in a real-world scenario")

workflow = [
    ("1", "Woman Downloads App", "Opens HormoCare AI on her smartphone browser. PWA installs to home screen instantly — no app store needed.", PINK),
    ("2", "Creates Health Profile", "Enters age (25), weight (62kg), height (160cm), selects family history of PCOS (mother had it). Takes 90 seconds.", TEAL),
    ("3", "Logs Daily Symptoms", "Each evening (< 1 min): marks period status, rates acne (3/5), selects hair fall level (High), logs mood (anxious), sleep (5 hrs).", ORANGE),
    ("4", "AI Analyses Patterns", "After 2 weeks of data, the algorithm detects: cycle > 35 days + high acne + weight gain pattern + family history → flags risk.", DARK_PINK),
    ("5", "Gets Risk Score", "Dashboard shows: 68% PCOS Risk (Yellow/Moderate). Lists 4 contributing factors. Shows trend graph of cycle irregularity over time.", RED),
    ("6", "Receives Care Plan", "Personalised output: 'Consult gynaecologist within 2 weeks' + Indian diet plan (low GI roti, dal, millets) + 30-min daily walk + Myo-inositol.", GREEN),
    ("7", "Shares Report with Doctor", "Generates a PDF with 2 weeks of symptom data, risk score, and trend graphs. Doctor gets instant context — no more 'it's just stress' dismissals.", BLUE),
]

for i, (num, title, desc, clr) in enumerate(workflow):
    y = Inches(1.4 + i * 0.8)
    # Number circle
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.6), y + Inches(0.05), Inches(0.5), Inches(0.5))
    circ.fill.solid(); circ.fill.fore_color.rgb = clr; circ.line.fill.background()
    txt(sl, Inches(0.6), y + Inches(0.08), Inches(0.5), Inches(0.42),
        num, size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Title
    txt(sl, Inches(1.3), y + Inches(0.02), Inches(2.5), Inches(0.35),
        title, size=13, bold=True, color=clr)
    # Description
    txt(sl, Inches(1.3), y + Inches(0.35), Inches(11), Inches(0.4),
        desc, size=11, color=DARK)


# ─── F-SLIDE 4: Competitive Analysis ─────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "Competitive Analysis", "How HormoCare AI compares to existing solutions in the market")

comp_data = [
    ("Feature", "Flo App", "Clue", "MyFitnessPal", "HormoCare AI"),
    ("PCOS Risk Detection", "✗ No", "✗ No", "✗ No", "✓ AI-Powered"),
    ("Multi-Symptom Tracking", "Period only", "Period + mood", "Weight only", "15+ parameters"),
    ("AI Risk Scoring", "✗ No", "✗ No", "✗ No", "✓ 0–100% score"),
    ("India-Specific Design", "✗ Global", "✗ Global", "✗ Global", "✓ Indian focus"),
    ("Personalised Diet Plan", "Generic tips", "✗ No", "Calorie only", "✓ PCOS Indian diet"),
    ("Doctor Report Export", "✗ No", "✗ No", "✗ No", "✓ PDF reports"),
    ("Works Offline", "✗ No", "✗ No", "✗ No", "✓ Full offline"),
    ("Privacy-First (No Cloud)", "Cloud-based", "Cloud-based", "Cloud-based", "✓ Local storage"),
    ("Supplement Guidance", "✗ No", "✗ No", "Generic", "✓ PCOS-specific"),
    ("Exercise Recommendations", "✗ No", "✗ No", "Generic", "✓ PCOS-tailored"),
]

rows = len(comp_data)
cols = 5
tbl_shape = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(1.5), Inches(12.1), Inches(5.0))
tbl = tbl_shape.table

col_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.3), Inches(3.0)]
for i, w in enumerate(col_widths):
    tbl.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = comp_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = "Arial"
            p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True
                p.font.color.rgb = WHITE
            elif c == 4:
                p.font.bold = True
                p.font.color.rgb = GREEN if "✓" in comp_data[r][c] else DARK
            elif "✗" in comp_data[r][c]:
                p.font.color.rgb = RED
            else:
                p.font.color.rgb = DARK
        if r == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_HDR
        elif c == 4:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF0, 0xFF, 0xF0)
        elif r % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = TABLE_ALT
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE


# ─── F-SLIDE 5: Real Impact (Non-Monetary) ───────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "The Change We Bring", "Impact measured in lives transformed — not revenue")

impacts = [
    ("🩺", "Slash Diagnosis Delay", "From 2–3 years to under 6 months. Early detection means a 25-year-old catches PCOS before it causes infertility at 30. That's a family saved.", RED),
    ("🔇", "Break the Silence", "Millions of Indian women suffer alone — told 'it's normal' by family. HormoCare shows them with data that their symptoms are real and treatable.", PINK),
    ("🌾", "Healthcare for Bharat", "In tier-2/3 cities and villages, specialist access is near zero. Our AI pre-screening bridges that gap — bringing intelligence where doctors can't reach.", TEAL),
    ("💪", "Data-Powered Self-Advocacy", "A woman walks into a clinic with a PDF showing 3 months of tracked symptoms and a 72% risk score. No doctor dismisses that as 'just stress' anymore.", ORANGE),
    ("🧠", "Mental Health Ripple Effect", "PCOS causes anxiety, depression, body image issues. Early intervention and support through our platform addresses the mental health crisis that silently follows.", DARK_PINK),
    ("👨‍👩‍👧", "Protect Entire Families", "Untreated PCOS → fertility struggles → relationship strain → financial burden of IVF. One early alert prevents years of suffering for the whole family.", GREEN),
    ("📚", "Health Literacy Revolution", "Every user learns about hormonal health, nutrition, and their own body. This knowledge passes to daughters, sisters, friends — creating generational awareness.", BLUE),
    ("🏥", "Reduce Healthcare System Load", "By catching PCOS early, we prevent the downstream surge of diabetes, cardiac issues, and fertility treatments that overburden India's hospitals.", MID_BLUE),
]

for i, (icon, title, desc, clr) in enumerate(impacts):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(1.4 + row * 1.4)
    card(sl, x, y, Inches(6.0), Inches(1.25), WHITE, clr)
    txt(sl, x + Inches(0.15), y + Inches(0.1), Inches(0.4), Inches(0.4),
        icon, size=18)
    txt(sl, x + Inches(0.6), y + Inches(0.08), Inches(5.1), Inches(0.3),
        title, size=13, bold=True, color=clr)
    txt(sl, x + Inches(0.6), y + Inches(0.4), Inches(5.1), Inches(0.75),
        desc, size=10, color=DARK)


# ─── F-SLIDE 6: Future Roadmap ───────────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "Future Roadmap", "What we build in the next 6–24 months")

phases = [
    ("Now → 3 Months", "Foundation", PINK, [
        "Finalise AI algorithm with gynaecologist input",
        "Beta test with 100 women across 3 cities",
        "Refine UI/UX based on user feedback",
        "Add Hindi language support",
        "Partner with 2 gynaecologist clinics for validation",
    ]),
    ("3 → 6 Months", "Scale & Integrate", TEAL, [
        "Launch native mobile app (React Native)",
        "Multilingual: Tamil, Bengali, Marathi, Telugu",
        "Doctor appointment booking integration",
        "Lab test partnership (Thyrocare, SRL)",
        "Community forum for peer support",
    ]),
    ("6 → 12 Months", "Platform Evolution", ORANGE, [
        "Upgrade to ML model (TensorFlow) with real data",
        "Wearable device integration (smartwatch data)",
        "Teleconsultation with PCOS specialists",
        "Pharmacy delivery for supplements",
        "ASHA worker outreach programme",
    ]),
    ("12 → 24 Months", "National & Beyond", GREEN, [
        "Integration with Ayushman Bharat Digital Mission",
        "Corporate wellness programme rollout",
        "Expand to thyroid, PCOD, endometriosis screening",
        "Southeast Asia & Middle East expansion",
        "Research partnerships with medical colleges",
    ]),
]

for i, (timeline, title, clr, items) in enumerate(phases):
    x = Inches(0.4 + i * 3.2)
    card(sl, x, Inches(1.5), Inches(3.0), Inches(5.2), LIGHT_BG, clr)
    s = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.5), Inches(3.0), Inches(0.65))
    s.fill.solid(); s.fill.fore_color.rgb = clr; s.line.fill.background()
    txt(sl, x, Inches(1.55), Inches(3.0), Inches(0.25),
        timeline, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(sl, x, Inches(1.8), Inches(3.0), Inches(0.3),
        title, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        txt(sl, x + Inches(0.2), Inches(2.3 + j * 0.55), Inches(2.6), Inches(0.5),
            f"▸  {item}", size=11, color=DARK)


# ─── F-SLIDE 7: Cost Breakdown & Sustainability ─────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)
stitle(sl, "Cost Breakdown & Sustainability Model")

# Development Cost table
cost_data = [
    ("Component", "Approach", "Cost Estimate", "Timeline"),
    ("Web App (MVP)", "HTML/CSS/JS + PWA", "₹0 (self-built)", "Completed"),
    ("Mobile App", "React Native", "₹1.5–2L (freelance) or self", "3 months"),
    ("Backend + DB", "Firebase (free tier)", "₹0–5K/month", "1 month"),
    ("AI Model Upgrade", "TensorFlow.js", "₹0 (open source)", "2 months"),
    ("Domain + Hosting", "Vercel / Netlify", "₹3–5K/year", "Immediate"),
    ("Medical Validation", "Gynaecologist consult", "₹20–30K", "Ongoing"),
    ("Total MVP Cost", "", "₹2–3 Lakhs", "6 months"),
]

rows = len(cost_data)
cols = 4
tbl_shape = sl.shapes.add_table(rows, cols, Inches(0.6), Inches(1.4), Inches(7), Inches(3.5))
tbl = tbl_shape.table

cw = [Inches(2.0), Inches(2.0), Inches(1.5), Inches(1.5)]
for i, w in enumerate(cw):
    tbl.columns[i].width = w

for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.text = cost_data[r][c]
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11); p.font.name = "Arial"; p.alignment = PP_ALIGN.CENTER
            if r == 0:
                p.font.bold = True; p.font.color.rgb = WHITE
            elif r == rows - 1:
                p.font.bold = True; p.font.color.rgb = DARK_PINK
            else:
                p.font.color.rgb = DARK
        if r == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_HDR
        elif r == rows - 1:
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF5)
        elif r % 2 == 0:
            cell.fill.solid(); cell.fill.fore_color.rgb = TABLE_ALT
        else:
            cell.fill.solid(); cell.fill.fore_color.rgb = WHITE

# Sustainability (non-monetary focus)
card(sl, Inches(8), Inches(1.4), Inches(4.7), Inches(5.2), LIGHT_BG, GREEN)
txt(sl, Inches(8.3), Inches(1.5), Inches(4.2), Inches(0.35),
    "♻️  Sustainability Model", size=16, bold=True, color=GREEN)

sust_pts = [
    ("Open-Source Core", "Free forever for basic risk assessment — ensuring no woman is excluded due to cost"),
    ("Government Partnership", "Align with National Digital Health Mission for funding and reach through public health channels"),
    ("Freemium for Advanced", "Premium features (detailed reports, teleconsult) sustain development costs"),
    ("NGO & CSR Funding", "Partner with women's health NGOs and corporate CSR programmes for subsidised access"),
    ("Academic Collaboration", "Medical college partnerships for research validation and continuous improvement"),
    ("Community-Driven", "User community contributes to awareness, reduces acquisition costs organically"),
]
for i, (title, desc) in enumerate(sust_pts):
    y = Inches(2.0 + i * 0.73)
    txt(sl, Inches(8.3), y, Inches(4.2), Inches(0.25),
        f"▸  {title}", size=11, bold=True, color=DARK)
    txt(sl, Inches(8.5), y + Inches(0.25), Inches(4.0), Inches(0.42),
        desc, size=10, color=GRAY)


# ─── F-SLIDE 8: Thank You ────────────────────────────────────────────
sl = prs2.slides.add_slide(prs2.slide_layouts[6])
set_bg(sl, 0xFF, 0xFF, 0xFF)
top_bar2(sl); bottom_bar2(sl); footer2(sl)

# Decorative circles
circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(0.5), Inches(5), Inches(5))
circ.fill.solid(); circ.fill.fore_color.rgb = RGBColor(0xFF, 0xF0, 0xF5); circ.line.fill.background()
circ2 = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(-2), Inches(3.5), Inches(4), Inches(4))
circ2.fill.solid(); circ2.fill.fore_color.rgb = RGBColor(0xF0, 0xF8, 0xFF); circ2.line.fill.background()

txt(sl, Inches(1.5), Inches(1.0), Inches(10), Inches(0.6),
    "🌸", size=48, align=PP_ALIGN.CENTER)
txt(sl, Inches(1.5), Inches(1.7), Inches(10), Inches(0.8),
    "HormoCare AI", size=44, bold=True, color=PINK, align=PP_ALIGN.CENTER)
accent_line(sl, Inches(5.5), Inches(2.55), Inches(2.3))

rich(sl, Inches(2), Inches(2.8), Inches(9.3), Inches(2.0), [
    {"t": '"Every woman deserves to know her body.', "s": 20, "c": DARK, "a": PP_ALIGN.CENTER, "sa": 4},
    {"t": 'Every symptom deserves to be heard.', "s": 20, "c": DARK, "a": PP_ALIGN.CENTER, "sa": 4},
    {"t": 'Every diagnosis deserves to come early."', "s": 20, "c": DARK, "a": PP_ALIGN.CENTER, "sa": 16},
    {"t": "— Team Sparsh", "s": 16, "b": True, "c": PINK, "a": PP_ALIGN.CENTER, "sa": 0},
])

txt(sl, Inches(2), Inches(4.8), Inches(9), Inches(0.4),
    "Anshika Madhu  •  Nilakshi  •  Anchal Kumari  •  Mayank Kumar  •  Arpan Gandhi  •  Saurabh Kumar",
    size=13, color=GRAY, align=PP_ALIGN.CENTER)

card(sl, Inches(4), Inches(5.5), Inches(5.3), Inches(1.0), LIGHT_BG, PINK)
txt(sl, Inches(4.2), Inches(5.6), Inches(4.9), Inches(0.35),
    "Demo Video:  [Insert 1-min demo link here]", size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)
txt(sl, Inches(4.2), Inches(5.95), Inches(4.9), Inches(0.3),
    "Thank You  |  Questions & Discussion", size=14, bold=True, color=PINK, align=PP_ALIGN.CENTER)


# ── Save Finale PPT ──────────────────────────────────────────────────
path2 = "/Users/mmt12368/Desktop/vtpaurelia/hormocare-ai/HormoCare_AI_SIH_Finale.pptx"
prs2.save(path2)
print(f"✓ Finale PPT saved: {path2}  ({len(prs2.slides)} slides)")
print("\nDone! Both presentations generated.")
